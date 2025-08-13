import ast
import random
import unicodedata
from pathlib import Path
from typing import List, Optional

import numpy as np
import pandas as pd
import typer
from loguru import logger
from pptx import Presentation

from src.config import (
    FLAGS_BY_CATEGORY,
    INTERIM_DATA_DIR,
    RAW_DATA_DIR,
    REFRIGERATED_CATEGORIES,
)
from src.utils.general_utils import convert_to_date

app = typer.Typer()


@app.command("transform")
def transform_data(
    input_file: str = "transactions.csv",
    sku_file: str = "sku.csv",
    output_file: str = "transaction_fpg.csv",
    item_id: str = "SDeptName",
):
    """Load and transform raw store data: rename columns, convert dates, normalize values."""
    df_transaction = pd.read_csv(RAW_DATA_DIR / input_file, encoding="utf-8-sig")
    df_sku = pd.read_csv(RAW_DATA_DIR / sku_file, encoding="utf-8-sig")
    df = df_transaction.merge(
        df_sku[["Sku", "Nh2", "Mfgr", "SDeptName"]], on="Sku", how="left"
    ).dropna()
    df = df[(df["Mfgr"].isin([2, 3, 4]))]
    df = df[["MergedId", item_id, "Sku", "Sales", "Quantity", "Date", "Cost"]].copy()
    df["Price"] = df["Sales"] / df["Quantity"]
    df = convert_to_date(df, "Date")
    df.rename(columns={item_id: "item_id"}, inplace=True)
    try:
        df.to_csv(
            INTERIM_DATA_DIR / output_file,
            index=False,
            date_format="%Y-%m-%d",
            encoding="utf-8-sig",
        )
        logger.success(f"Data transformed and saved to {output_file}")
        return df
    except Exception as e:
        logger.error(f"Error processing key {df['Key'][0]}: {e}")
        return None, None


from src.config import (
    EXTERNAL_DATA_DIR,
    INTERIM_DATA_DIR,
    PROCESSED_DATA_DIR,
    RAW_DATA_DIR,
)


@app.command("extract-layout")
def extract_layout(
    layout_pptx: str = typer.Option(
        "layout.pptx", "--pptx", "-p", help="Tên file PPTX trong EXTERNAL_DATA_DIR"
    ),
    output_csv: str = typer.Option(
        "layout.csv", "--out", "-o", help="Tên file CSV sẽ lưu vào INTERIM_DATA_DIR"
    ),
    cold: List[str] = typer.Option(
        REFRIGERATED_CATEGORIES, "--cold", help="Danh mục lạnh, khai báo nhiều lần"
    ),
):
    """
    Đọc shapes từ PPTX (mm) và xuất CSV gồm:
    Category, x, y, width, height, is_refrigerated, is_entrance, is_cashier
    """
    # 1) Chuẩn hoá chuỗi
    norm = lambda s: "".join(
        ch
        for ch in unicodedata.normalize("NFD", (s or "").strip().lower())
        if unicodedata.category(ch) != "Mn"
    )
    # 2) Chuẩn hoá tập lạnh
    cold_set = {norm(c) for c in (cold or [])}
    norm_flags = {
        flag: {norm(cat) for cat in (cats or [])}
        for flag, cats in FLAGS_BY_CATEGORY.items()
    }

    # 4) Đọc PPTX
    EMU_PER_MM = 914400 / 25.4
    prs = Presentation(EXTERNAL_DATA_DIR / layout_pptx)
    rows = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                txt = (sh.text_frame.text or "").strip()
                if not txt:
                    continue
                n = norm(txt)
                row = {
                    "Category": txt,
                    "x": sh.left / EMU_PER_MM,
                    "y": sh.top / EMU_PER_MM,
                    "width": sh.width / EMU_PER_MM,
                    "height": sh.height / EMU_PER_MM,
                    "is_refrigerated": int(n in cold_set),
                }
                for flag in norm_flags.keys():
                    row[flag] = int(n in norm_flags[flag])
                rows.append(row)

    if not rows:
        raise RuntimeError("No shapes with text found in PPTX.")

    df = pd.DataFrame(rows).drop_duplicates(keep="first")
    # 6) Đảm bảo hai cột tối thiểu tồn tại ngay cả khi FLAGS_BY_CATEGORY không khai báo
    for must in ["is_entrance", "is_cashier"]:
        if must not in df.columns:
            df[must] = 0
    out_path = INTERIM_DATA_DIR / output_csv
    df.to_csv(out_path, index=False)
    typer.echo(f"Wrote {len(df)} rows to {out_path}")


class DataLoader:
    """
    Đọc & chuẩn hoá dữ liệu cho tối ưu layout.
    - Loại Entrance/Cashier khỏi slots ngay từ đầu.
    - Cache bảng slots đã sort theo (y, x) để dùng lại nhanh.
    - Xây dựng all_items từ rules + itemsets + categories của slot bày hàng.
    """

    def __init__(
        self,
        assoc_rules_path,
        freq_itemsets_path,
        layout_real_path,
        margin_matrix_path: Optional[str] = None,
    ):
        self.assoc_rules = pd.read_csv(assoc_rules_path)
        self.freq_itemsets = pd.read_csv(freq_itemsets_path)
        self.layout_real = pd.read_csv(layout_real_path).drop_duplicates(keep="first")
        self.margin_matrix = (
            pd.read_csv(margin_matrix_path, index_col=0)
            if margin_matrix_path is not None
            else None
        )
        self._slots_xy = None  # cache
        self.positions = None
        self.all_items: List[str] = []
        self.refrig_cats: List[str] = []
        self._process()

    def _process(self):
        # Bảo đảm cột mặc định tồn tại
        for col in ["is_refrigerated", "is_entrance", "is_cashier", "width", "height"]:
            if col not in self.layout_real.columns:
                self.layout_real[col] = 0

        # Vị trí: yêu cầu có x,y
        if {"x", "y"}.issubset(self.layout_real.columns):
            self.positions = list(zip(self.layout_real["x"], self.layout_real["y"]))
        else:
            raise ValueError("layout_real.csv thiếu cột x,y")

        # Parse rules / itemsets
        antecedents = self.assoc_rules["antecedent"].apply(ast.literal_eval)
        consequents = self.assoc_rules["consequent"].apply(ast.literal_eval)
        itemsets = self.freq_itemsets["items"].apply(ast.literal_eval)

        all_items = set()
        for ser in antecedents.tolist() + consequents.tolist():
            all_items.update(ser)
        for sublist in itemsets.tolist():
            all_items.update(sublist)

        # Lọc SLOT BÀY HÀNG (loại Entrance/Cashier) & cache
        df = self.layout_real.copy()
        mask_slots = (df["is_entrance"].fillna(0).astype(int) != 1) & (
            df["is_cashier"].fillna(0).astype(int) != 1
        )
        slots = (
            df.loc[
                mask_slots, ["Category", "x", "y", "width", "height", "is_refrigerated"]
            ]
            .assign(
                width=lambda d: d["width"].fillna(0),
                height=lambda d: d["height"].fillna(0),
            )
            .sort_values(["y", "x"])
            .reset_index(drop=True)
        )
        self._slots_xy = slots[["Category", "x", "y", "width", "height"]].copy()

        # Cập nhật all_items chỉ với categories của slot bày hàng
        layout_cats = slots["Category"].dropna().astype(str).tolist()
        all_items.update(layout_cats)
        self.all_items = sorted(all_items)

        # Refrigerated cats chỉ trong slot bày hàng
        self.refrig_cats = (
            slots.loc[slots["is_refrigerated"].fillna(0).astype(int) == 1, "Category"]
            .astype(str)
            .tolist()
        )

    # helpers
    def sorted_slots_xy(self) -> pd.DataFrame:
        # trả bản cache (đã sort & fill width/height)
        return self._slots_xy.copy()


if __name__ == "__main__":
    app()
